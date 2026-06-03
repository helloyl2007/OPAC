import os
import subprocess
import logging
import asyncio
from fastapi import APIRouter, UploadFile, HTTPException, Body, Depends, Request
from fastapi.responses import JSONResponse
from pathlib import Path
import uuid
import time
import fitz  # PyMuPDF
from pptx import Presentation
from pydantic import BaseModel
from datetime import datetime
from typing import List, Optional

# 导入相关模块
from app.utils.video_generator import generate_video
import json
from sqlalchemy.orm import Session
from app.core.database import get_db, SessionLocal
from app.models.models import VideoFile, PPTFile
from app.api.deps import get_current_user
from app.utils.task_manager import task_manager
from app.api.video_utils import save_video_record, update_video_status, VIDEO_STATUS

router = APIRouter()

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 确保目录存在
os.makedirs("static/ppt_upload", exist_ok=True)
os.makedirs("static/ppt_upload/pdf_tmp", exist_ok=True)
os.makedirs("static/ppt_upload/ppt_img_tmp", exist_ok=True)


async def convert_to_images(pdf_path: str, output_dir: str) -> list:
    """将PDF转换为图片"""
    try:
        # 打开PDF文件
        doc = fitz.open(pdf_path)
        images = []
        
        # 遍历每一页
        for page_num in range(len(doc)):
            # 获取页面
            page = doc[page_num]
            
            # 获取页面原始尺寸
            rect = page.rect
            
            # 计算缩放比例，使高度为720px
            target_height = 720
            zoom = target_height / rect.height
            
            # 设置渲染参数
            mat = fitz.Matrix(zoom, zoom)
            
            # 将页面渲染为图片
            pix = page.get_pixmap(matrix=mat)
            img_path = os.path.join(output_dir, f"slide_{page_num + 1:03d}.png")
            pix.save(img_path)
            images.append(f"slide_{page_num + 1:03d}.png")
        
        doc.close()
        return images
        
    except Exception as e:
        logger.error(f"PDF转图片失败: {str(e)}")
        raise

@router.post("/upload")
async def upload_ppt(file: UploadFile):
    try:
        # 生成唯一标识符
        unique_id = f"{int(time.time())}_{uuid.uuid4().hex[:8]}"
        original_name = Path(file.filename).stem
        
        # 构建文件路径
        ppt_path = f"static/ppt_upload/{unique_id}.pptx"
        pdf_path = f"static/ppt_upload/pdf_tmp/{unique_id}.pdf"
        img_dir = f"static/ppt_upload/ppt_img_tmp/{unique_id}_slides"
        
        # 创建图片目录
        os.makedirs(img_dir, exist_ok=True)

        # 保存上传的PPT文件
        with open(ppt_path, "wb") as f:
            content = await file.read()
            f.write(content)
        
        # 保存原始文件名信息到JSON文件，确保后续处理可以获取
        info_path = f"static/ppt_upload/{unique_id}_info.json"
        with open(info_path, "w", encoding="utf-8") as f:
            json.dump({
                "original_filename": file.filename,
                "upload_time": time.time(),
                "file_size": len(content)
            }, f, ensure_ascii=False)
            
        # 提取PPT备注
        prs = Presentation(ppt_path)
        notes = []
        for slide in prs.slides:
            note = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
            notes.append(note)
        
        # 使用libreoffice转换PPT为PDF
        convert_cmd = [
            'libreoffice',
            '--headless',
            '--convert-to',
            'pdf',
            '--outdir',
            os.path.dirname(pdf_path),
            ppt_path
        ]
        
        logger.info(f"Converting PPT to PDF: {original_name}")
        process = subprocess.Popen(convert_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        stdout, stderr = process.communicate()
        
        if process.returncode != 0:
            raise HTTPException(status_code=500, detail=f"PPT转PDF失败: {stderr.decode()}")
            
        # 将PDF转换为图片
        logger.info("Converting PDF to images")
        images = await convert_to_images(pdf_path, img_dir)
        
        if not images:
            raise HTTPException(status_code=500, detail="未能生成任何图片")

        logger.info(f"Generated {len(images)} slides")
        
        # 返回结果时包含备注
        return {
            "id": unique_id,
            "original_filename": file.filename,
            "status": "success",
            "total_slides": len(images),
            "images": [f"/static/ppt_upload/ppt_img_tmp/{unique_id}_slides/{img}" for img in images],
            "notes": notes  # 添加备注数组
        }
        
    except Exception as e:
        logger.error(f"处理失败: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# 添加请求模型
class SystemPPTRequest(BaseModel):
    pptId: str

@router.post("/convert-system-ppt")
async def convert_system_ppt(
    request: SystemPPTRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """转换系统已有的PPT文件"""
    try:
        logger.info(f"Converting system PPT, id: {request.pptId}")
        
        # 修改这里：从数据库获取PPT文件路径，而不是直接构建
        ppt_file = db.query(PPTFile).filter(PPTFile.id == request.pptId).first()
        if not ppt_file:
            logger.error(f"PPT record not found in database: {request.pptId}")
            raise HTTPException(status_code=404, detail="PPT文件记录不存在")
            
        # 使用数据库中的file_path，确保调整到正确的绝对路径
        if os.path.isabs(ppt_file.file_path):
            ppt_path = ppt_file.file_path
        else:
            # 假设相对路径以static开头
            if ppt_file.file_path.startswith("static/"):
                ppt_path = ppt_file.file_path
            else:
                ppt_path = os.path.join("static", ppt_file.file_path)
            
            # 确保使用正确的路径分隔符
            ppt_path = ppt_path.replace("/", os.sep)
        
        if not os.path.exists(ppt_path):
            logger.error(f"PPT file not found at path: {ppt_path}")
            raise HTTPException(status_code=404, detail="PPT文件不存在")
            
        # 生成唯一标识符
        unique_id = f"{int(time.time())}_{uuid.uuid4().hex[:8]}"
        
        # 保存原始文件名信息到JSON文件
        # 修改：使用PPT的title作为视频名称，而不是文件名
        original_filename = ppt_file.title if ppt_file.title else os.path.basename(ppt_path)
        info_path = f"static/ppt_upload/{unique_id}_info.json"
        with open(info_path, "w", encoding="utf-8") as f:
            json.dump({
                "original_filename": original_filename,
                "ppt_title": ppt_file.title,  # 添加title字段
                "source": "system",
                "ppt_id": request.pptId
            }, f, ensure_ascii=False)
        
        # 构建输出路径，保持与上传PPT相同的命名规则
        pdf_name = f"{unique_id}.pdf"  # PDF文件名
        pdf_path = os.path.join("static", "ppt_upload", "pdf_tmp", pdf_name)
        img_dir = os.path.join("static", "ppt_upload", "ppt_img_tmp", f"{unique_id}_slides")
        
        # 确保目录存在
        os.makedirs(img_dir, exist_ok=True)
        os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
        
        logger.info(f"Converting PPT to PDF: {ppt_path} -> {pdf_path}")
        
        # 提取PPT备注
        try:
            prs = Presentation(ppt_path)
            notes = []
            for slide in prs.slides:
                note = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
                notes.append(note)
        except Exception as e:
            logger.error(f"提取PPT备注失败: {str(e)}")
            notes = []

        # 转换为PDF
        try:
            convert_cmd = [
                'libreoffice',
                '--headless',
                '--convert-to',
                'pdf',
                '--outdir',
                os.path.abspath(os.path.dirname(pdf_path)),
                os.path.abspath(ppt_path)
            ]
            
            logger.info(f"Running command: {' '.join(convert_cmd)}")
            process = subprocess.Popen(convert_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            stdout, stderr = process.communicate()
            
            if process.returncode != 0:
                logger.error(f"PDF conversion failed: {stderr.decode()}")
                raise HTTPException(status_code=500, detail=f"PPT转PDF失败: {stderr.decode()}")
        except Exception as e:
            logger.error(f"执行转换命令失败: {str(e)}")
            raise HTTPException(status_code=500, detail=f"执行转换命令失败: {str(e)}")

        # 重命名生成的PDF文件以匹配期望的命名规则
        try:
            generated_pdf = os.path.join(os.path.dirname(pdf_path), 
                                       os.path.splitext(os.path.basename(ppt_path))[0] + '.pdf')
            if os.path.exists(generated_pdf):
                os.rename(generated_pdf, pdf_path)
                logger.info(f"重命名PDF文件: {generated_pdf} -> {pdf_path}")
            else:
                logger.warning(f"生成的PDF文件不存在: {generated_pdf}")
        except Exception as e:
            logger.error(f"重命名PDF文件失败: {str(e)}")
            # 继续执行，因为可能文件名已经正确
            
        # 等待文件生成完成
        await asyncio.sleep(1)
        
        # 验证PDF文件是否存在
        if not os.path.exists(pdf_path):
            logger.error(f"PDF file not found after conversion: {pdf_path}")
            raise HTTPException(status_code=500, detail="PDF文件转换失败，请检查后台日志")
            
        # 转换为图片
        try:
            images = await convert_to_images(pdf_path, img_dir)
            
            if not images:
                logger.error("未能生成任何图片")
                raise HTTPException(status_code=500, detail="未能生成任何图片")
        except Exception as e:
            logger.error(f"PDF转图片失败: {str(e)}")
            raise HTTPException(status_code=500, detail=f"PDF转图片失败: {str(e)}")
            
        # 构建响应
        response_data = {
            "id": unique_id,
            "original_filename": os.path.basename(ppt_path),
            "status": "success",
            "total_slides": len(images),
            "images": [f"/static/ppt_upload/ppt_img_tmp/{unique_id}_slides/{img}" for img in images],
            "notes": notes
        }
        
        logger.info(f"系统PPT转换成功，生成 {len(images)} 张图片")
        return response_data
        
    except HTTPException:
        # 直接重抛HTTP异常
        raise
    except Exception as e:
        logger.error(f"转换系统PPT失败: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

# 添加视频生成请求模型
class VideoGenerationRequest(BaseModel):
    sessionId: str
    notes: List[str]
    duration: int = 5  # 默认每张幻灯片显示5秒
    voiceType: str = "zh-CN-YunyangNeural"  # 默认语音类型
    video_id: Optional[int] = None  # 如果是编辑现有视频，提供视频ID

# 视频生成请求响应模型
class VideoGenerationResponse(BaseModel):
    task_id: str
    status: str
    message: str
    video_id: int

# 新增：获取视频编辑详情API
@router.get("/edit/{video_id}")
async def get_video_for_edit(
    video_id: int,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """获取视频编辑信息"""
    try:
        # 检查是否为管理员
        is_admin = "admin" in current_user.roles if current_user.roles else False
        
        # 构建查询
        query = db.query(VideoFile).filter(VideoFile.id == video_id)
        
        # 如果不是管理员，只能编辑自己的视频
        if not is_admin:
            query = query.filter(VideoFile.user_id == current_user.id)
            
        video = query.first()
        
        if not video:
            raise HTTPException(status_code=404, detail="视频不存在或无权编辑")
        
        # 从session_id获取原始的PPT信息
        session_id = video.session_id
        if not session_id:
            raise HTTPException(status_code=400, detail="无法编辑此视频，缺少会话信息")
            
        # 构建图片目录路径
        img_dir = os.path.join("static", "ppt_upload", "ppt_img_tmp", f"{session_id}_slides")
        
        # 检查图片目录是否存在
        if not os.path.exists(img_dir):
            raise HTTPException(status_code=404, detail="未找到原始幻灯片")
            
        # 获取图片列表并排序
        image_files = [f for f in os.listdir(img_dir) if f.endswith(('.png', '.jpg', '.jpeg'))]
        image_files.sort()  # 确保顺序正确
        
        if not image_files:
            raise HTTPException(status_code=404, detail="未找到任何图片")
        
        # 在响应中发送更详细的日志用于调试
        logger.info(f"获取视频编辑信息 - 视频ID: {video_id}, 会话ID: {session_id}")
        if video.meta_info:
            logger.info(f"视频meta_info: {video.meta_info[:200]}...")
        else:
            logger.info("视频meta_info为空")
        
        # 获取之前的解说文本
        notes = []
        metadata = None
        
        # 尝试从meta_info中获取解说文本
        if video.meta_info and video.meta_info.strip():
            try:
                metadata = json.loads(video.meta_info)
                if 'notes' in metadata:
                    notes = metadata.get('notes', [])
                    logger.info(f"从meta_info中获取到解说文本数量: {len(notes)}")
                else:
                    logger.warning("meta_info中不包含notes字段")
            except Exception as e:
                logger.error(f"解析meta_info失败: {str(e)}")
                
        # 如果没有获取到解说文本，就创建空列表
        if not notes:
            logger.warning(f"没有找到解说文本，为每张幻灯片创建空笔记，图片数量: {len(image_files)}")
            notes = [''] * len(image_files)
        
        # 确保notes长度和图片数量一致
        if len(notes) < len(image_files):
            logger.warning(f"解说文本数量({len(notes)})少于图片数量({len(image_files)})，补充空文本")
            notes.extend([''] * (len(image_files) - len(notes)))
        elif len(notes) > len(image_files):
            logger.warning(f"解说文本数量({len(notes)})多于图片数量({len(image_files)})，截断多余文本")
            notes = notes[:len(image_files)]
                
        # 构建响应
        result = {
            "id": video.id,
            "session_id": session_id,
            "filename": video.filename,
            "images": [f"/static/ppt_upload/ppt_img_tmp/{session_id}_slides/{img}" for img in image_files],
            "notes": notes,
            "voice_type": metadata.get('voice_type', 'aixia') if metadata else 'aixia',
            "duration": metadata.get('duration', 3) if metadata else 3
        }
        
        logger.info(f"返回视频编辑信息: {result}")
        return result
    
    except Exception as e:
        logger.error(f"获取视频编辑信息失败: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

# 修改视频生成API
@router.post("/generate-video")
async def generate_video_api(
    request: VideoGenerationRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """生成视频API - 使用后台任务异步处理"""
    try:
        session_id = request.sessionId
        user_id = current_user.id  # 获取用户ID
        username = current_user.username  # 记录用户名，便于调试
        
        # 记录详细的请求信息
        logger.info(f"视频生成请求 - 用户: {username}(ID:{user_id}), 会话: {session_id}")
        
        # 构建图片目录路径
        img_dir = os.path.join("static", "ppt_upload", "ppt_img_tmp", f"{session_id}_slides")
        
        # 检查图片目录是否存在
        if not os.path.exists(img_dir):
            raise HTTPException(status_code=404, detail=f"未找到会话ID: {session_id} 的图片")
        
        # 获取图片列表并排序
        image_files = [f for f in os.listdir(img_dir) if f.endswith(('.png', '.jpg', '.jpeg'))]
        image_files.sort()  # 确保顺序正确
        
        if not image_files:
            raise HTTPException(status_code=404, detail="未找到任何图片")
        
        # 构建图片路径列表
        image_paths = [os.path.join(img_dir, img) for img in image_files]
        
        # 确保视频输出目录存在
        os.makedirs("static/generated/videos", exist_ok=True)
        
        # 新增：检查是否是编辑已有视频
        video_id = None
        video_path = None
        
        if request.video_id:
            # 验证权限
            query = db.query(VideoFile).filter(VideoFile.id == request.video_id)
            if not "admin" in current_user.roles:
                query = query.filter(VideoFile.user_id == current_user.id)
                
            existing_video = query.first()
            if existing_video:
                video_id = existing_video.id
                video_path = existing_video.filepath  # 使用原始文件路径
                logger.info(f"编辑现有视频 - ID: {video_id}, 路径: {video_path}")
            else:
                raise HTTPException(status_code=404, detail="视频不存在或无权编辑")
        else:
            # 生成唯一的视频文件名
            video_filename = f"{session_id}_{int(time.time())}.mp4"
            video_path = os.path.join("static", "generated", "videos", video_filename)
            # 确保路径使用正斜杠
            video_path = video_path.replace("\\", "/")
            
        # 获取文件名信息
        original_filename = None
        ppt_title = None
        info_path = f"static/ppt_upload/{session_id}_info.json"
        if os.path.exists(info_path):
            try:
                with open(info_path, "r", encoding="utf-8") as f:
                    info_data = json.load(f)
                    # 优先使用ppt_title字段
                    if "ppt_title" in info_data:
                        ppt_title = info_data["ppt_title"]
                    if "original_filename" in info_data:
                        original_filename = Path(info_data["original_filename"]).stem
            except Exception as e:
                logger.warning(f"读取文件信息失败: {str(e)}")
                
        # 如果没有获取到原始文件名或标题，使用默认名称
        video_title = ppt_title or original_filename
        if not video_title:
            video_title = f"视频_{int(time.time())}"
            
        # 保存当前的生成参数到meta_info
        meta_info = {
            'notes': request.notes,
            'voice_type': request.voiceType,
            'duration': request.duration,
            'progress': 0
        }
        
        logger.info(f"视频生成参数 - notes长度: {len(request.notes)}, 语音类型: {request.voiceType}, 每页时长: {request.duration}")
        
        # 创建或更新视频记录
        if request.video_id:
            # 获取现有视频的meta_info，确保不丢失信息
            existing_video = db.query(VideoFile).filter(VideoFile.id == request.video_id).first()
            if existing_video:
                logger.info(f"更新现有视频记录 - ID: {existing_video.id}, 原meta_info存在: {existing_video.meta_info is not None}")
                
                # 更新视频记录，设置状态为处理中
                existing_video.status = VIDEO_STATUS["PROCESSING"]
                
                # 保存新的元数据，包含解说文本
                logger.info(f"保存新的meta_info，包含解说文本数量: {len(request.notes)}")
                existing_video.meta_info = json.dumps(meta_info)
                
                # 提交更改
                db.commit()
                
                video_id = existing_video.id
                video_path = existing_video.filepath
                
                # 验证meta_info是否更新成功
                updated_video = db.query(VideoFile).filter(VideoFile.id == video_id).first()
                if updated_video:
                    try:
                        updated_meta = json.loads(updated_video.meta_info)
                        logger.info(f"验证meta_info - notes长度: {len(updated_meta.get('notes', []))}")
                    except Exception as e:
                        logger.error(f"验证meta_info失败: {str(e)}")
            else:
                raise HTTPException(status_code=404, detail="视频不存在或无权编辑")
        else:
            # 创建新的视频记录，状态为处理中
            video_id = save_video_record(
                filepath=video_path,
                user_id=user_id,
                session_id=session_id,
                original_filename=video_title,  # 使用title作为视频名称
                status=VIDEO_STATUS["PROCESSING"]
            )
            
            # 单独更新meta_info字段
            try:
                db_video = db.query(VideoFile).filter(VideoFile.id == video_id).first()
                if db_video:
                    logger.info(f"为新视频设置meta_info，包含解说文本数量: {len(request.notes)}")
                    db_video.meta_info = json.dumps(meta_info)
                    db.commit()
            except Exception as e:
                logger.error(f"更新新视频meta_info失败: {str(e)}")

        # 定义异步生成视频任务
        async def process_video_generation():
            # 显式使用变量，确保值被正确传入闭包
            _user_id = user_id
            _session_id = session_id
            _video_id = video_id
            _video_path = video_path  # 确保使用正确的视频路径
            
            logger.info(f"视频生成任务开始 - 用户ID: {_user_id}, 视频ID: {_video_id}")
            
            try:
                # 设置初始进度，确保不覆盖notes
                update_status_result = update_video_status(_video_id, VIDEO_STATUS["PROCESSING"], progress=5)
                logger.info(f"更新初始进度结果: {update_status_result}")
                
                # 准备处理图片和音频
                total_steps = len(image_paths) * 2 + 1  # 图片+音频+合成
                current_step = 0
                
                # 修改视频生成函数调用，添加进度回调
                async def progress_callback(step, total, message):
                    nonlocal current_step
                    current_step = step
                    progress = min(95, (step / total) * 100)  # 最大95%，留5%给最终处理
                    update_video_status(_video_id, VIDEO_STATUS["PROCESSING"], progress=progress)
                
                # 调用视频生成工具，添加进度回调参数
                result = await generate_video(
                    image_paths=image_paths,
                    output_path=_video_path,
                    notes=request.notes,
                    slide_duration=request.duration,
                    voice_type=request.voiceType,
                    progress_callback=progress_callback
                )
                
                if not result["success"]:
                    logger.error(f"视频生成失败: {result['error']}")
                    update_video_status(_video_id, VIDEO_STATUS["FAILED"], error=result['error'])
                    return {"status": "failed", "error": result['error']}
                
                # 更新视频记录，保持notes信息
                db_session = SessionLocal()
                try:
                    # 先读取当前meta_info，确保保留notes
                    video = db_session.query(VideoFile).filter(VideoFile.id == _video_id).first()
                    if video and video.meta_info:
                        try:
                            current_meta = json.loads(video.meta_info)
                            # 确保保留原有的notes
                            if 'notes' not in current_meta and request.notes:
                                logger.info("meta_info中没有notes字段，重新添加")
                                current_meta['notes'] = request.notes
                                
                            # 更新状态和进度
                            current_meta['progress'] = 100
                            video.meta_info = json.dumps(current_meta)
                            video.status = VIDEO_STATUS["COMPLETED"]
                            video.duration = int(result["duration"])
                            db_session.commit()
                            logger.info("成功更新视频完成状态并保留notes")
                        except Exception as e:
                            logger.error(f"读取/更新meta_info失败: {str(e)}")
                finally:
                    db_session.close()
                
                logger.info(f"视频生成完成 - 用户ID: {_user_id}, 视频ID: {_video_id}")
                
                # 返回视频信息
                video_url_path = f"/static/generated/videos/{os.path.basename(_video_path)}"
                # 确保URL路径使用正斜杠
                video_url_path = video_url_path.replace("\\", "/")
                return {
                    "status": "success",
                    "message": "视频生成成功",
                    "video_url": video_url_path,
                    "video_path": _video_path.replace("\\", "/"),  # 确保使用正斜杠
                    "duration": result["duration"],
                    "video_id": _video_id,
                    "user_id": _user_id
                }
                
            except Exception as e:
                logger.error(f"视频生成过程中出错: {str(e)}", exc_info=True)
                update_video_status(_video_id, VIDEO_STATUS["FAILED"], error=str(e))
                return {"status": "failed", "error": str(e)}
        
        # 记录任务创建信息
        logger.info(f"创建视频后台任务 - 用户ID: {user_id}")
        
        # 创建后台任务
        task_id = task_manager.create_task(process_video_generation)
        
        # 记录任务ID
        logger.info(f"任务已创建 - ID: {task_id}, 用户ID: {user_id}")
        
        # 立即返回任务ID和视频ID
        return {
            "task_id": task_id,
            "video_id": video_id,
            "status": "processing",
            "message": "视频生成任务已启动，请稍后检查状态",
            "user_id": user_id
        }
        
    except Exception as e:
        logger.error(f"启动视频生成失败: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

# 添加保存草稿的请求模型
class SaveDraftRequest(BaseModel):
    sessionId: str
    notes: List[str]
    duration: int = 5
    voiceType: str = "zh-CN-YunyangNeural" 
    video_id: Optional[int] = None

@router.post("/save-draft")
async def save_video_draft(
    request: SaveDraftRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """保存视频草稿，简化版本"""
    try:
        # 基本参数验证
        if not request.sessionId or not request.notes:
            return {"status": "error", "message": "缺少必要参数"}
        
        session_id = request.sessionId
        user_id = current_user.id
        
        # 简单记录
        logger.info(f"保存草稿 - 用户:{user_id}, 会话:{session_id}, notes数量:{len(request.notes)}")
        
        # 构建基础元数据
        meta_info = {
            'notes': request.notes,
            'voice_type': request.voiceType,
            'duration': request.duration
        }
        
        # 如果提供了video_id，更新现有记录
        if request.video_id:
            try:
                # 直接查询视频，简化权限验证
                video = db.query(VideoFile).filter(
                    VideoFile.id == request.video_id,
                    VideoFile.user_id == user_id
                ).first()
                
                if not video:
                    # 查询失败，创建新记录
                    logger.warning(f"找不到要更新的视频ID: {request.video_id}，将创建新记录")
                    raise ValueError("找不到要更新的视频，将创建新记录")
                    
                # 更新元数据
                video.meta_info = json.dumps(meta_info)
                db.commit()
                
                return {
                    "status": "success", 
                    "message": "更新成功",
                    "video_id": video.id
                }
            except Exception as e:
                logger.warning(f"更新失败，将创建新记录: {str(e)}")
                # 失败后继续创建新记录
        
        # 创建新记录前，先检查是否已有相同session_id的记录
        existing_video = db.query(VideoFile).filter(
            VideoFile.session_id == session_id,
            VideoFile.user_id == user_id
        ).first()
        
        if existing_video:
            # 如果找到了相同session_id的记录，更新它而不是创建新的
            logger.info(f"发现相同session_id的记录，更新而非创建: {existing_video.id}")
            existing_video.meta_info = json.dumps(meta_info)
            db.commit()
            
            return {
                "status": "success",
                "message": "更新现有记录成功",
                "video_id": existing_video.id
            }
        
        # 创建新记录流程        
        # 生成临时文件路径
        temp_path = os.path.join("static", "generated", "videos", f"draft_{session_id}.mp4")
        # 确保路径使用正斜杠
        temp_path = temp_path.replace("\\", "/")
        
        # 确保目录存在
        os.makedirs(os.path.dirname(temp_path), exist_ok=True)
        
        # 获取原始文件名信息
        original_filename = None
        info_path = f"static/ppt_upload/{session_id}_info.json"
        if os.path.exists(info_path):
            try:
                with open(info_path, "r", encoding="utf-8") as f:
                    info_data = json.load(f)
                    if "original_filename" in info_data:
                        original_filename = Path(info_data["original_filename"]).stem
            except Exception as e:
                logger.warning(f"读取信息文件失败: {str(e)}")
                
        # 如果没有获取到原始文件名，使用默认名称
        if not original_filename:
            original_filename = f"草稿_{session_id}"
        
        # 创建视频记录
        new_video = VideoFile(
            filename=original_filename,
            filepath=temp_path,
            user_id=user_id,
            status=VIDEO_STATUS["PENDING"],
            session_id=session_id,
            meta_info=json.dumps(meta_info)
        )
        
        db.add(new_video)
        db.commit()
        db.refresh(new_video)
        
        logger.info(f"创建新视频记录 - ID: {new_video.id}")
        
        return {
            "status": "success",
            "message": "保存成功", 
            "video_id": new_video.id
        }
        
    except Exception as e:
        db.rollback()
        logger.error(f"保存草稿失败: {str(e)}", exc_info=True)
        # 返回用户友好的错误信息
        return {"status": "error", "message": f"保存失败: {str(e)}"}

# 添加任务状态查询接口
@router.get("/task/{task_id}")
async def get_task_status(task_id: str):
    """查询任务状态"""
    status = task_manager.get_task_status(task_id)
    if status["status"] == "not_found":
        raise HTTPException(status_code=404, detail="任务不存在")
    return status

@router.get("/video-status/{session_id}")
async def get_video_status(session_id: str):
    """获取视频生成状态"""
    try:
        # 检查是否存在对应的视频
        video_dir = os.path.join("static", "generated", "videos")
        video_files = [f for f in os.listdir(video_dir) if f.startswith(session_id)]
        
        if not video_files:
            return {"status": "not_found"}
            
        # 返回最新生成的视频
        latest_video = sorted(video_files)[-1]
        video_path = os.path.join(video_dir, latest_video)
        
        return {
            "status": "completed",
            "video_url": f"/static/generated/videos/{latest_video}",
            "file_size": os.path.getsize(video_path)
        }
        
    except Exception as e:
        logger.error(f"获取视频状态失败: {str(e)}")
        return {"status": "error", "message": str(e)}