from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.utils.logger import logger
import copy
import io

def recreate_slide_by_pptx(src_path: str, save_path: str, indexs: list):
    """使用python-pptx重新创建幻灯片"""
    src_prs = Presentation(src_path)
    dst_prs = Presentation(src_path)
    
    # 清除目标演示文稿中的所有现有幻灯片
    for _ in range(len(dst_prs.slides._sldIdLst)):
        rId = dst_prs.slides._sldIdLst[0].rId
        dst_prs.part.drop_rel(rId)
        dst_prs.slides._sldIdLst.remove(dst_prs.slides._sldIdLst[0])

    # 按照指定顺序复制幻灯片
    for idx in indexs:
        src_slide = src_prs.slides[idx]
        dst_slide = dst_prs.slides.add_slide(src_slide.slide_layout)
        
        # 处理所有形状
        for shape in src_slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 获取原始图片的所有变换信息
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    rotation = shape.rotation if hasattr(shape, 'rotation') else 0
                    
                    # 创建图片流
                    image_stream = io.BytesIO(shape.image.blob)
                    
                    # 添加图片并保持原有属性
                    pic = dst_slide.shapes.add_picture(
                        image_stream,
                        left,
                        top,
                        width,
                        height
                    )
                    
                    # 设置旋转角度
                    pic.rotation = rotation
                    
                    # 复制变换属性
                    if hasattr(shape.element, 'xfrm'):
                        src_xfrm = shape.element.xfrm
                        dst_xfrm = pic.element.xfrm
                        for attr in ['flipH', 'flipV', 'rot']:
                            if hasattr(src_xfrm, attr):
                                setattr(dst_xfrm, attr, getattr(src_xfrm, attr))
                else:
                    # 复制非图片形状
                    el = shape.element
                    newel = copy.deepcopy(el)
                    dst_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
            except Exception as e:
                logger.error(f"Error copying shape: {str(e)}")
                continue

    # 保存新文件
    dst_prs.save(save_path)
