from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.utils.logger import logger
from app.services.ppt_tools import recreate_slide_by_pptx
import json
import re
import os
import random

class PPTGenerator:
    PPT_PARAM_PATTERN = r'\{(.*?)\}'
    
    def __init__(self, template_path=None):
        try:
            self.prs = Presentation(template_path) if template_path else Presentation()
            self.template_path = template_path
            self.template_params = {}
            
            if template_path:
                self._extract_template_params()
            else:
                self._setup_default_layouts()
        except Exception as e:
            logger.error(f"初始化PPT生成器失败: {str(e)}")
            raise

    def _setup_default_layouts(self):
        """设置默认布局"""
        # ... 保持原有代码 ...
        pass
    
    def _extract_template_params(self):
        """提取模板中的参数"""
        try:
            # 定义页面索引
            self.template_params = {
                "first_slide": {"nos": [0], "params": []},
                "catalogue_slide": {"nos": [1], "params": []},
                "title_slide": {"nos": [2], "params": []},
                "content_slide": {"nos": [3, 4, 5, 6, 7, 8], "params": []},
                "end_slide": {"nos": [9], "params": []}
            }

            # 提取每种类型页面的参数
            for slide_type, info in self.template_params.items():
                nos = info["nos"]
                for n in nos:
                    if n < len(self.prs.slides):
                        slide = self.prs.slides[n]
                        params = {}
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame"):
                                matches = re.findall(self.PPT_PARAM_PATTERN, shape.text_frame.text)
                                if matches:
                                    params[shape.name] = matches
                        info["params"].append(params)
            
            logger.info(f"提取到的模板参数: {self.template_params}")
        except Exception as e:
            logger.error(f"提取模板参数失败: {str(e)}")
            raise

    def _process_content_params(self, page_content):
        """处理内容页参数，匹配模板占位符"""
        params = {}
        
        # 添加标题和子标题
        params['title'] = page_content.get('title', '')
        sub_title = page_content.get('sub_title', '')
        params['sub_title'] = sub_title
        params['subtitle'] = sub_title  # 兼容 subtitle 格式
        params['sub_title_1'] = sub_title  # 兼容编号形式
        
        # 处理描述和内容点，保持一一对应
        desc_list = page_content.get('desc', [])
        contents = page_content.get('contents', [])
        
        # 确保desc和contents长度一致
        max_len = max(len(desc_list), len(contents))
        desc_list = desc_list + [''] * (max_len - len(desc_list))
        contents = contents + [''] * (max_len - len(contents))
        
        # 按序号添加对应的描述和内容
        for idx, (desc, content) in enumerate(zip(desc_list, contents), start=1):
            params[f'desc_{idx}'] = desc
            params[f'content_{idx}'] = content
            params[f'sub_title_{idx}'] = sub_title  # 每个内容点都带有子标题
        
        logger.info(f"Processed content params: {params}")
        return params

    def _add_slide_notes(self, slide, params: dict):
        """为幻灯片添加备注，保持内容顺序"""
        notes_slide = slide.notes_slide
        notes_text = []
        
        # 1. 处理首页和目录页
        if 'topic' in params:
            notes_text.append(f"{params['topic']}")
        # 处理目录页的标题列表
        title_keys = sorted([k for k in params.keys() if k.startswith('title_')])
        for key in title_keys:
            if params[key]:
                notes_text.append(params[key])
                
        # 2. 处理标题和子标题
        if 'title' in params and params['title']:
            notes_text.append(f"{params['title']}")
        if 'sub_title' in params and params['sub_title']:
            notes_text.append(f"{params['sub_title']}")
            
        # 3. 按顺序添加内容
        idx = 1
        while f'desc_{idx}' in params or f'content_{idx}' in params:
            # 先添加描述
            if f'desc_{idx}' in params and params[f'desc_{idx}']:
                notes_text.append(params[f'desc_{idx}'])
            # 再添加内容
            if f'content_{idx}' in params and params[f'content_{idx}']:
                notes_text.append(params[f'content_{idx}'])
            idx += 1
        
        # 写入备注
        if notes_text:
            notes_slide.notes_text_frame.text = "\n".join(notes_text)

    def generate_ppt(self, topic: str, content: dict, output_path: str):
        """生成PPT"""
        try:
            # 解析JSON内容
            if isinstance(content, str):
                content = json.loads(content)
            
            # 处理页面内容
            pages = content.get("pages", [])
            if not pages:
                raise ValueError("No pages found in content")

            # 准备页面参数和顺序
            all_params = []
            final_nos = []
            
            # 1. 首页参数
            all_params.append({"topic": topic})
            final_nos.append(self.template_params["first_slide"]["nos"][0])
            
            # 2. 目录页参数
            catalogue_params = {f'title_{i+1}': page["title"] for i, page in enumerate(pages)}
            all_params.append(catalogue_params)
            final_nos.append(self.template_params["catalogue_slide"]["nos"][0])
            
            # 3. 处理内容页，每个章节的每个子标题生成一个内容页
            content_layouts = self.template_params["content_slide"]["nos"]
            # 创建一个可用模板队列，用于不重复随机选择
            available_layouts = list(content_layouts)
            
            for page in pages:
                # 添加章节标题页
                all_params.append({"title": page["title"]})
                final_nos.append(self.template_params["title_slide"]["nos"][0])
                
                # 为每个子页面（子标题）生成一个内容页
                for subpage in page.get("pages", []):
                    content_params = self._process_content_params(subpage)
                    all_params.append(content_params)
                    
                    # 不重复随机选择模板
                    if not available_layouts:
                        # 如果所有模板都已使用，重新填充可用模板列表
                        available_layouts = list(content_layouts)
                    
                    # 从可用模板中选择一个
                    chosen_layout = random.choice(available_layouts)
                    # 从可用列表中移除已选择的模板
                    available_layouts.remove(chosen_layout)
                    
                    final_nos.append(chosen_layout)
            
            # 4. 结束页
            all_params.append({"topic": topic})  # 可以在结束页显示主题
            final_nos.append(self.template_params["end_slide"]["nos"][0])

            # 使用ppt_tools重新创建幻灯片
            recreate_slide_by_pptx(self.template_path, output_path, final_nos)
            
            # 填充内容
            ppt = Presentation(output_path)
            for idx, params in enumerate(all_params):
                slide = ppt.slides[idx]
                # 处理形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text = run.text
                                matches = re.finditer(self.PPT_PARAM_PATTERN, text)
                                for match in matches:
                                    param_name = match.group(1)
                                    placeholder = "{" + param_name + "}"
                                    if param_name in params:
                                        text = text.replace(placeholder, str(params[param_name]))
                                    else:
                                        text = text.replace(placeholder, "")
                                run.text = text
                
                # 为每页添加备注，特殊处理结束页
                if idx == len(all_params) - 1:  # 最后一页（结束页）
                    slide.notes_slide.notes_text_frame.text = "同学们都记住了吗？我们下次见哦！"
                else:
                    self._add_slide_notes(slide, params)
            
            # 保存最终的PPT
            ppt.save(output_path)
            
        except Exception as e:
            logger.error(f"生成PPT失败: {str(e)}", exc_info=True)
            raise

    def _generate_title_slide(self, title: str):
        """生成标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_placeholder = slide.shapes.title
        if title_placeholder:
            title_placeholder.text = title
        else:
            # 如果没有标题占位符，创建文本框
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(1.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = title
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _generate_content_slide(self, page_content: dict):
        """生成内容页"""
        # 获取随机布局索引
        layout_indices = self.template_params["content_slide"]["index"]
        layout_index = random.choice(layout_indices)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_index])
        
        # 对每个形状处理占位符
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            
            text_frame = shape.text_frame
            original_text = text_frame.text
            
            # 如果没有占位符，跳过处理
            if '{' not in original_text:
                continue
            
            # 清除现有文本准备重新填充
            text_frame.clear()
            
            # 处理不同类型的占位符
            if '{title}' in original_text:
                text_frame.text = page_content.get('title', '')
            elif '{subtitle}' in original_text:
                text_frame.text = page_content.get('sub_title', '')
            elif '{desc}' in original_text or '{content}' in original_text:
                # 如果是内容区域，添加带格式的内容
                for subpage in page_content.get('pages', []):
                    if '{desc}' in original_text and subpage.get('desc'):
                        p = text_frame.add_paragraph()
                        p.text = '• ' + subpage['desc']
                        p.level = 0
                        
                    if '{content}' in original_text and subpage.get('content'):
                        p = text_frame.add_paragraph()
                        p.text = '  ' + subpage['content']  # 使用缩进表示层级
                        p.level = 1
            
            # 设置段落格式
            for paragraph in text_frame.paragraphs:
                if paragraph.level == 0:
                    paragraph.font.size = Pt(24)
                elif paragraph.level == 1:
                    paragraph.font.size = Pt(18)
                paragraph.font.name = "微软雅黑"
                paragraph.space_after = Pt(12)

    def save(self, output_path: str):
        """保存PPT文件"""
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.prs.save(output_path)
        except Exception as e:
            logger.error(f"保存PPT失败: {str(e)}")
            raise
